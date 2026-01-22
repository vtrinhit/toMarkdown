"""Custom converter with enhanced image/shape extraction.

This converter provides best-in-class handling for:
- Excel files: Extracts embedded images, renders shapes/charts as images via LibreOffice
- PDF files: Extracts all embedded images with correct positioning
- PowerPoint files: Extracts slide images and embedded media

All images are embedded as Base64 in the markdown output for portability.
"""

import asyncio
import base64
import subprocess
import tempfile
from pathlib import Path
from typing import List, Optional

from .base import BaseConverter


class CustomConverter(BaseConverter):
    """
    Custom converter with enhanced image extraction.

    Provides superior handling of embedded images and shapes
    compared to other converters. Uses LibreOffice for rendering
    charts and vector shapes.
    """

    name = "Custom (Enhanced)"
    description = "Enhanced converter with full image/shape/chart/hyperlink extraction. Best for Office files with images, PDF, and PowerPoint."
    supported_extensions: List[str] = [
        "xlsx", "xls", "xlsm", "xlsb",  # Excel
        "docx", "doc",                   # Word
        "pdf",                           # PDF
        "pptx", "ppt",                   # PowerPoint
    ]

    async def convert(self, file_path: Path, output_path: Path) -> str:
        """Convert file to markdown with embedded images."""
        ext = file_path.suffix.lower().lstrip(".")

        if ext in ["xlsx", "xls", "xlsm", "xlsb"]:
            return await self._convert_excel(file_path, output_path)
        elif ext in ["docx", "doc"]:
            return await self._convert_word(file_path, output_path)
        elif ext == "pdf":
            return await self._convert_pdf(file_path, output_path)
        elif ext in ["pptx", "ppt"]:
            return await self._convert_pptx(file_path, output_path)
        else:
            raise ValueError(f"Unsupported file type: {ext}")

    async def _convert_excel(self, file_path: Path, output_path: Path) -> str:
        """Convert Excel file with images, shapes, and charts rendered correctly."""

        def _process():
            from openpyxl import load_workbook

            wb = load_workbook(file_path, data_only=True)
            markdown_parts = []
            markdown_parts.append(f"# {file_path.stem}\n")

            # Check if any sheet has charts - if so, render via LibreOffice
            has_charts = False
            chart_positions = {}  # {sheet_idx: [(chart_row, chart_info), ...]}
            for sheet_idx, sheet_name in enumerate(wb.sheetnames):
                ws = wb[sheet_name]
                if hasattr(ws, '_charts') and ws._charts:
                    has_charts = True
                    chart_positions[sheet_idx] = self._get_chart_positions(ws)

            # Render sheets with LibreOffice if there are charts
            rendered_sheets = {}
            if has_charts:
                rendered_sheets = self._render_excel_with_libreoffice(file_path)

            for sheet_idx, sheet_name in enumerate(wb.sheetnames):
                ws = wb[sheet_name]
                markdown_parts.append(f"\n## Sheet: {sheet_name}\n")

                # Extract images with their Excel row positions
                images_by_excel_row = self._extract_excel_images_with_positions(ws)

                # Add rendered chart images to images_by_excel_row
                if sheet_idx in chart_positions and sheet_idx in rendered_sheets:
                    for chart_row, chart_title in chart_positions[sheet_idx]:
                        # The rendered sheet contains the chart
                        # We'll add the full rendered image at the chart position
                        if chart_row not in images_by_excel_row:
                            images_by_excel_row[chart_row] = []
                        images_by_excel_row[chart_row].append(
                            (rendered_sheets[sheet_idx], "png", None, f"Chart: {chart_title}")
                        )

                # Extract cell data with original Excel row numbers and hyperlinks
                row_data = []  # [(excel_row, cells), ...]
                hyperlinks = self._extract_excel_hyperlinks(ws)

                for excel_row_idx, row in enumerate(ws.iter_rows(values_only=True), start=1):
                    # Skip completely empty rows
                    if any(cell is not None for cell in row):
                        cells = []
                        for col_idx, cell in enumerate(row, start=1):
                            cell_value = str(cell) if cell is not None else ""
                            # Check if this cell has a hyperlink
                            cell_key = (excel_row_idx, col_idx)
                            if cell_key in hyperlinks:
                                link = hyperlinks[cell_key]
                                # Format as markdown link
                                display_text = cell_value if cell_value else link
                                cell_value = f"[{display_text}]({link})"
                            cells.append(cell_value)
                        row_data.append((excel_row_idx, cells))

                if row_data:
                    # Build markdown with images inline
                    markdown_parts.append(
                        self._build_excel_markdown_with_images(row_data, images_by_excel_row)
                    )

                # Any images that couldn't be positioned (no anchor) go at the end
                orphan_images = images_by_excel_row.get(-1, [])
                if orphan_images:
                    for idx, img_tuple in enumerate(orphan_images, 1):
                        img_data = img_tuple[0]
                        img_format = img_tuple[1]
                        label = img_tuple[3] if len(img_tuple) > 3 else f"Image {idx}"
                        base64_img = base64.b64encode(img_data).decode('utf-8')
                        mime_type = f"image/{img_format.lower()}"
                        if img_format.lower() == "jpg":
                            mime_type = "image/jpeg"
                        markdown_parts.append(
                            f"\n![{label}](data:{mime_type};base64,{base64_img})\n"
                        )

            wb.close()

            markdown_content = "\n".join(markdown_parts)
            output_file = output_path / f"{file_path.stem}.md"
            output_file.write_text(markdown_content, encoding="utf-8")
            return str(output_file)

        loop = asyncio.get_running_loop()
        return await loop.run_in_executor(None, _process)

    async def _convert_word(self, file_path: Path, output_path: Path) -> str:
        """Convert Word document with images, shapes, and hyperlinks."""

        def _process():
            from docx import Document
            from docx.oxml.ns import qn

            doc = Document(file_path)
            markdown_parts = []
            markdown_parts.append(f"# {file_path.stem}\n")

            # Extract images from document
            images = self._extract_word_images(doc)
            image_index = 0

            # Process paragraphs and tables
            for element in doc.element.body:
                if element.tag.endswith('p'):  # Paragraph
                    para = None
                    for p in doc.paragraphs:
                        if p._element is element:
                            para = p
                            break

                    if para:
                        # Check for images in paragraph
                        has_image = False
                        for run in para.runs:
                            if run._element.xpath('.//a:blip'):
                                has_image = True
                                break

                        if has_image and image_index < len(images):
                            img_data, img_format = images[image_index]
                            base64_img = base64.b64encode(img_data).decode('utf-8')
                            markdown_parts.append(
                                f"\n![Image {image_index + 1}](data:image/{img_format};base64,{base64_img})\n"
                            )
                            image_index += 1

                        # Extract text with hyperlinks
                        text = self._extract_word_paragraph_with_links(para)
                        if text.strip():
                            # Handle headings
                            if para.style and para.style.name.startswith('Heading'):
                                level = 1
                                try:
                                    level = int(para.style.name.replace('Heading ', ''))
                                except ValueError:
                                    level = 1
                                markdown_parts.append(f"{'#' * level} {text}\n")
                            else:
                                markdown_parts.append(text + "\n")

                elif element.tag.endswith('tbl'):  # Table
                    for table in doc.tables:
                        if table._element is element:
                            markdown_parts.append(self._convert_word_table(table) + "\n")
                            break

            markdown_content = "\n".join(markdown_parts)
            output_file = output_path / f"{file_path.stem}.md"
            output_file.write_text(markdown_content, encoding="utf-8")
            return str(output_file)

        loop = asyncio.get_running_loop()
        return await loop.run_in_executor(None, _process)

    def _extract_word_images(self, doc) -> list:
        """Extract images from Word document.

        Returns:
            list: [(img_data, img_format), ...]
        """
        images = []
        seen_hashes = set()

        try:
            # Method 1: Extract from document relationships
            for rel in doc.part.rels.values():
                if "image" in rel.reltype:
                    try:
                        img_data = rel.target_part.blob
                        img_hash = hash(img_data)
                        if img_hash in seen_hashes:
                            continue
                        seen_hashes.add(img_hash)

                        # Determine format from content type
                        content_type = rel.target_part.content_type
                        if "png" in content_type:
                            img_format = "png"
                        elif "jpeg" in content_type or "jpg" in content_type:
                            img_format = "jpeg"
                        elif "gif" in content_type:
                            img_format = "gif"
                        else:
                            img_format = "png"

                        images.append((img_data, img_format))
                    except Exception:
                        continue
        except Exception:
            pass

        return images

    def _extract_word_paragraph_with_links(self, para) -> str:
        """Extract paragraph text with hyperlinks preserved."""
        from docx.oxml.ns import qn

        result_parts = []

        for child in para._element:
            # Handle hyperlinks
            if child.tag == qn('w:hyperlink'):
                # Get the relationship ID
                r_id = child.get(qn('r:id'))
                link_url = None

                if r_id:
                    try:
                        rel = para.part.rels.get(r_id)
                        if rel:
                            link_url = rel.target_ref
                    except Exception:
                        pass

                # Get text from hyperlink
                link_text = ""
                for run_elem in child.findall(qn('w:r')):
                    for text_elem in run_elem.findall(qn('w:t')):
                        if text_elem.text:
                            link_text += text_elem.text

                if link_url and link_text:
                    result_parts.append(f"[{link_text}]({link_url})")
                elif link_text:
                    result_parts.append(link_text)

            # Handle regular runs
            elif child.tag == qn('w:r'):
                for text_elem in child.findall(qn('w:t')):
                    if text_elem.text:
                        result_parts.append(text_elem.text)

        return "".join(result_parts)

    def _convert_word_table(self, table) -> str:
        """Convert Word table to markdown format."""
        rows = []

        for row in table.rows:
            cells = []
            for cell in row.cells:
                # Get cell text, escape pipes
                cell_text = cell.text.strip().replace("|", "\\|").replace("\n", " ")
                cells.append(cell_text)
            rows.append(cells)

        if not rows:
            return ""

        # Build markdown table
        lines = []

        # Normalize column count
        max_cols = max(len(row) for row in rows)
        rows = [row + [""] * (max_cols - len(row)) for row in rows]

        # Header
        lines.append("| " + " | ".join(rows[0]) + " |")
        lines.append("| " + " | ".join(["---"] * max_cols) + " |")

        # Data rows
        for row in rows[1:]:
            lines.append("| " + " | ".join(row) + " |")

        return "\n".join(lines)

    def _render_excel_with_libreoffice(self, file_path: Path) -> dict:
        """Render Excel sheets as images using LibreOffice headless.

        Returns:
            dict: {sheet_index: image_bytes, ...}
        """
        sheet_images = {}

        try:
            with tempfile.TemporaryDirectory() as tmp_dir:
                tmp_path = Path(tmp_dir)

                # Try different LibreOffice executable names
                soffice_cmd = self._find_libreoffice_executable()
                if not soffice_cmd:
                    return sheet_images

                # Convert Excel to PDF using LibreOffice
                result = subprocess.run(
                    [
                        soffice_cmd, '--headless', '--convert-to', 'pdf',
                        '--outdir', str(tmp_path),
                        str(file_path)
                    ],
                    capture_output=True,
                    timeout=120,
                    check=False
                )

                pdf_file = tmp_path / f"{file_path.stem}.pdf"
                if pdf_file.exists():
                    import fitz  # PyMuPDF

                    doc = fitz.open(pdf_file)
                    for page_idx, page in enumerate(doc):
                        # Render page at 150 DPI for good quality
                        matrix = fitz.Matrix(150 / 72, 150 / 72)
                        pix = page.get_pixmap(matrix=matrix)
                        sheet_images[page_idx] = pix.tobytes("png")
                    doc.close()
        except subprocess.TimeoutExpired:
            pass
        except Exception:
            pass

        return sheet_images

    def _find_libreoffice_executable(self) -> Optional[str]:
        """Find LibreOffice executable on the system."""
        import shutil
        import platform

        # Common executable names
        candidates = ['soffice', 'libreoffice', 'libreoffice7.6', 'libreoffice7.5']

        # Windows-specific paths
        if platform.system() == 'Windows':
            candidates.extend([
                r'C:\Program Files\LibreOffice\program\soffice.exe',
                r'C:\Program Files (x86)\LibreOffice\program\soffice.exe',
            ])

        for cmd in candidates:
            if shutil.which(cmd):
                return cmd
            # Check if it's a direct path
            if Path(cmd).exists():
                return cmd

        return None

    def _get_chart_positions(self, ws) -> list:
        """Get chart positions (row numbers) from worksheet.

        Returns:
            list: [(row_number, chart_title), ...]
        """
        chart_positions = []

        try:
            for chart in ws._charts:
                # Get chart anchor position
                row = -1
                title = ""

                if hasattr(chart, 'anchor'):
                    anchor = chart.anchor
                    if hasattr(anchor, '_from'):
                        row = anchor._from.row + 1
                    elif hasattr(anchor, 'row'):
                        row = anchor.row + 1

                if hasattr(chart, 'title') and chart.title:
                    title = str(chart.title)
                else:
                    title = type(chart).__name__

                chart_positions.append((row, title))
        except Exception:
            pass

        return chart_positions

    def _extract_excel_hyperlinks(self, ws) -> dict:
        """Extract hyperlinks from worksheet.

        Returns:
            dict: {(row, col): url, ...}
        """
        hyperlinks = {}

        try:
            # Method 1: Get hyperlinks from worksheet._hyperlinks
            if hasattr(ws, '_hyperlinks'):
                for link in ws._hyperlinks:
                    try:
                        if hasattr(link, 'ref') and hasattr(link, 'target'):
                            # Parse cell reference (e.g., "A1", "B2")
                            from openpyxl.utils import coordinate_to_tuple
                            row, col = coordinate_to_tuple(link.ref.split(':')[0])
                            if link.target:
                                hyperlinks[(row, col)] = link.target
                    except Exception:
                        continue

            # Method 2: Iterate through cells to find hyperlinks
            for row in ws.iter_rows():
                for cell in row:
                    try:
                        if cell.hyperlink and cell.hyperlink.target:
                            hyperlinks[(cell.row, cell.column)] = cell.hyperlink.target
                    except Exception:
                        continue
        except Exception:
            pass

        return hyperlinks

    def _extract_sheet_data(self, ws) -> List[List[str]]:
        """Extract cell data from worksheet."""
        data = []
        for row in ws.iter_rows(values_only=True):
            # Skip completely empty rows
            if any(cell is not None for cell in row):
                data.append([str(cell) if cell is not None else "" for cell in row])
        return data

    def _build_excel_markdown_with_images(
        self, row_data: list, images_by_excel_row: dict
    ) -> str:
        """Build markdown table with images inserted at correct positions.

        Args:
            row_data: [(excel_row_num, [cell1, cell2, ...]), ...]
            images_by_excel_row: {excel_row: [(img_data, img_format, col), ...], ...}
        """
        if not row_data:
            return ""

        # Find all rows (including image-only rows)
        all_excel_rows = set(excel_row for excel_row, _ in row_data)
        for excel_row in images_by_excel_row:
            if excel_row > 0:  # Skip orphan images (-1)
                all_excel_rows.add(excel_row)

        # Sort all row numbers
        sorted_rows = sorted(all_excel_rows)

        # Build data lookup
        row_data_lookup = {excel_row: cells for excel_row, cells in row_data}

        # Determine max columns
        max_cols = max(len(cells) for _, cells in row_data) if row_data else 1

        lines = []
        table_started = False

        for excel_row in sorted_rows:
            # Check if this row has data
            has_data = excel_row in row_data_lookup
            # Check if this row has images
            has_images = excel_row in images_by_excel_row

            if has_data:
                cells = row_data_lookup[excel_row]
                # Normalize to max_cols
                cells = cells + [""] * (max_cols - len(cells))
                # Escape pipe characters
                escaped_cells = [cell.replace("|", "\\|") for cell in cells]

                if not table_started:
                    # Start table with header
                    lines.append("| " + " | ".join(escaped_cells) + " |")
                    lines.append("| " + " | ".join(["---"] * len(escaped_cells)) + " |")
                    table_started = True
                else:
                    lines.append("| " + " | ".join(escaped_cells) + " |")

            if has_images:
                # Close table before images if table is started
                if table_started:
                    lines.append("")  # Empty line to close table

                # Add images
                for img_idx, img_tuple in enumerate(
                    images_by_excel_row[excel_row], 1
                ):
                    img_data = img_tuple[0]
                    img_format = img_tuple[1]
                    # Optional label (4th element)
                    label = img_tuple[3] if len(img_tuple) > 3 else f"Row {excel_row} Image {img_idx}"

                    base64_img = base64.b64encode(img_data).decode("utf-8")
                    mime_type = f"image/{img_format.lower()}"
                    if img_format.lower() == "jpg":
                        mime_type = "image/jpeg"
                    lines.append(
                        f"![{label}](data:{mime_type};base64,{base64_img})"
                    )

                lines.append("")  # Empty line after images

                # Reset table state - next data row will start new table
                table_started = False

        return "\n".join(lines) + "\n"

    def _extract_excel_images_with_positions(self, ws) -> dict:
        """Extract embedded images from worksheet with their row positions.

        Returns:
            dict: {row_number: [(img_data, img_format, col), ...], ...}
                  row_number -1 is used for images without position
        """
        images_by_row = {}
        seen_hashes = set()  # Avoid duplicates

        try:
            for image in ws._images:
                try:
                    # Get image data
                    img_data = image._data()

                    # Skip duplicates based on content hash
                    img_hash = hash(img_data)
                    if img_hash in seen_hashes:
                        continue
                    seen_hashes.add(img_hash)

                    # Determine format
                    img_format = "png"
                    if hasattr(image, "format"):
                        img_format = image.format or "png"
                    elif hasattr(image, "_path") and image._path:
                        img_format = Path(image._path).suffix.lstrip(".") or "png"

                    # Get anchor position (row number)
                    row = -1  # Default: no position
                    col = None
                    if hasattr(image, "anchor"):
                        anchor = image.anchor
                        # TwoCellAnchor or OneCellAnchor
                        if hasattr(anchor, "_from"):
                            row = anchor._from.row + 1  # 0-indexed to 1-indexed
                            col = anchor._from.col + 1
                        elif hasattr(anchor, "row"):
                            row = anchor.row + 1
                            col = getattr(anchor, "col", None)
                            if col is not None:
                                col += 1

                    if row not in images_by_row:
                        images_by_row[row] = []
                    images_by_row[row].append((img_data, img_format, col))
                except Exception:
                    continue
        except Exception:
            pass

        return images_by_row

    async def _convert_pdf(self, file_path: Path, output_path: Path) -> str:
        """Convert PDF with images and hyperlinks positioned correctly relative to text."""

        def _process():
            import fitz  # PyMuPDF

            doc = fitz.open(file_path)
            markdown_parts = []
            markdown_parts.append(f"# {file_path.stem}\n")

            for page_num in range(len(doc)):
                page = doc[page_num]
                markdown_parts.append(f"\n## Page {page_num + 1}\n")

                # Extract hyperlinks from page
                page_links = self._extract_pdf_links(page)

                # Get all content blocks with positions
                content_items = []

                # Extract text blocks with positions
                text_blocks = page.get_text("blocks")
                for block in text_blocks:
                    if block[6] == 0:  # Text block (not image)
                        y0 = block[1]  # Top y-coordinate
                        x0, x1, y1 = block[0], block[2], block[3]
                        text = block[4].strip()
                        if text:
                            # Apply hyperlinks to text block
                            text_with_links = self._apply_pdf_links_to_text(
                                text, (x0, y0, x1, y1), page_links
                            )
                            content_items.append(("text", text_with_links, y0))

                # Extract images with positions
                seen_hashes = set()
                for img in page.get_images():
                    try:
                        xref = img[0]
                        base_image = page.parent.extract_image(xref)

                        if base_image:
                            img_data = base_image["image"]

                            # Skip duplicates
                            img_hash = hash(img_data)
                            if img_hash in seen_hashes:
                                continue
                            seen_hashes.add(img_hash)

                            # Skip very small images (icons/bullets)
                            if len(img_data) <= 1000:
                                continue

                            img_ext = base_image.get("ext", "png")

                            # Get image position on page
                            img_rects = page.get_image_rects(xref)
                            y_pos = 0
                            if img_rects:
                                y_pos = img_rects[0].y0  # Top of image

                            content_items.append(("image", (img_data, img_ext), y_pos))
                    except Exception:
                        continue

                # Sort by vertical position (top to bottom)
                content_items.sort(key=lambda x: x[2])

                # Output in order
                img_counter = 0
                for content_type, content, _ in content_items:
                    if content_type == "text":
                        # Clean and add text
                        cleaned = self._clean_text_block(content)
                        if cleaned:
                            markdown_parts.append(cleaned + "\n")
                    else:
                        img_counter += 1
                        img_data, img_format = content
                        base64_img = base64.b64encode(img_data).decode("utf-8")
                        markdown_parts.append(
                            f"\n![Page {page_num + 1} Image {img_counter}](data:image/{img_format};base64,{base64_img})\n"
                        )

            doc.close()

            markdown_content = "\n".join(markdown_parts)
            output_file = output_path / f"{file_path.stem}.md"
            output_file.write_text(markdown_content, encoding="utf-8")
            return str(output_file)

        loop = asyncio.get_running_loop()
        return await loop.run_in_executor(None, _process)

    def _extract_pdf_links(self, page) -> list:
        """Extract hyperlinks from a PDF page.

        Returns:
            list: [(rect, uri), ...] where rect is (x0, y0, x1, y1)
        """
        links = []

        try:
            for link in page.get_links():
                if link.get("kind") == 2:  # URI link
                    uri = link.get("uri", "")
                    rect = link.get("from")
                    if uri and rect:
                        links.append(((rect.x0, rect.y0, rect.x1, rect.y1), uri))
        except Exception:
            pass

        return links

    def _apply_pdf_links_to_text(self, text: str, block_rect: tuple, page_links: list) -> str:
        """Apply hyperlinks to text if the block intersects with any link.

        Args:
            text: The text content
            block_rect: (x0, y0, x1, y1) of the text block
            page_links: [(rect, uri), ...] from _extract_pdf_links

        Returns:
            Text with markdown links applied
        """
        bx0, by0, bx1, by1 = block_rect

        for link_rect, uri in page_links:
            lx0, ly0, lx1, ly1 = link_rect

            # Check if rectangles intersect
            if not (bx1 < lx0 or bx0 > lx1 or by1 < ly0 or by0 > ly1):
                # This text block is within a link area
                # Wrap the entire text in a markdown link
                return f"[{text}]({uri})"

        return text

    def _clean_text_block(self, text: str) -> str:
        """Clean a single text block from PDF."""
        lines = text.split("\n")
        cleaned = [line.strip() for line in lines if line.strip()]
        return " ".join(cleaned) if cleaned else ""

    async def _convert_pptx(self, file_path: Path, output_path: Path) -> str:
        """Convert PowerPoint with images and hyperlinks."""

        def _process():
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE

            prs = Presentation(file_path)
            markdown_parts = []
            markdown_parts.append(f"# {file_path.stem}\n")

            # Track all image hashes to avoid duplicates
            seen_hashes = set()

            for slide_num, slide in enumerate(prs.slides, 1):
                markdown_parts.append(f"\n## Slide {slide_num}\n")

                # Extract text and images together, respecting shape order
                content_parts = []
                slide_images = []

                for shape in slide.shapes:
                    try:
                        # Handle picture shapes
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            image = shape.image
                            img_data = image.blob
                            img_hash = hash(img_data)

                            if img_hash not in seen_hashes:
                                seen_hashes.add(img_hash)
                                img_format = image.ext.lower()
                                if img_format in ["jpg", "jpeg"]:
                                    img_format = "jpeg"

                                # Check if image has hyperlink
                                img_link = self._get_pptx_shape_hyperlink(shape)
                                slide_images.append((img_data, img_format, shape.top, img_link))

                        # Handle text shapes
                        elif hasattr(shape, "text") and shape.text.strip():
                            text = shape.text.strip()
                            # Check for hyperlinks in text frame
                            text_with_links = self._extract_pptx_text_with_links(shape)
                            content_parts.append((text_with_links, shape.top))
                    except Exception:
                        continue

                # Sort by vertical position and output
                all_content = []
                for text, pos in content_parts:
                    all_content.append(("text", text, pos))
                for img_tuple in slide_images:
                    img_data, img_format, pos = img_tuple[0], img_tuple[1], img_tuple[2]
                    img_link = img_tuple[3] if len(img_tuple) > 3 else None
                    all_content.append(("image", (img_data, img_format, img_link), pos))

                # Sort by position (top)
                all_content.sort(key=lambda x: x[2] if x[2] is not None else 0)

                img_counter = 0
                for content_type, content, _ in all_content:
                    if content_type == "text":
                        markdown_parts.append(content + "\n")
                    else:
                        img_counter += 1
                        img_data, img_format, img_link = content
                        base64_img = base64.b64encode(img_data).decode("utf-8")
                        img_md = f"![Slide {slide_num} Image {img_counter}](data:image/{img_format};base64,{base64_img})"
                        if img_link:
                            img_md = f"[{img_md}]({img_link})"
                        markdown_parts.append(f"\n{img_md}\n")

            markdown_content = "\n".join(markdown_parts)
            output_file = output_path / f"{file_path.stem}.md"
            output_file.write_text(markdown_content, encoding="utf-8")
            return str(output_file)

        loop = asyncio.get_running_loop()
        return await loop.run_in_executor(None, _process)

    def _get_pptx_shape_hyperlink(self, shape) -> Optional[str]:
        """Get hyperlink from a PowerPoint shape."""
        try:
            if hasattr(shape, "click_action") and shape.click_action:
                if hasattr(shape.click_action, "hyperlink") and shape.click_action.hyperlink:
                    return shape.click_action.hyperlink.address
        except Exception:
            pass
        return None

    def _extract_pptx_text_with_links(self, shape) -> str:
        """Extract text from PowerPoint shape with hyperlinks preserved."""
        try:
            if not hasattr(shape, "text_frame"):
                return shape.text.strip() if hasattr(shape, "text") else ""

            result_parts = []
            for paragraph in shape.text_frame.paragraphs:
                para_parts = []
                for run in paragraph.runs:
                    text = run.text
                    if not text:
                        continue

                    # Check if run has hyperlink
                    link = None
                    try:
                        if hasattr(run, "hyperlink") and run.hyperlink:
                            link = run.hyperlink.address
                    except Exception:
                        pass

                    if link:
                        para_parts.append(f"[{text}]({link})")
                    else:
                        para_parts.append(text)

                if para_parts:
                    result_parts.append("".join(para_parts))

            return "\n".join(result_parts) if result_parts else (shape.text.strip() if hasattr(shape, "text") else "")
        except Exception:
            return shape.text.strip() if hasattr(shape, "text") else ""
